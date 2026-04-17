"""
Script para criar apresentação PPTX da plataforma de reuniões - identidade visual SPOT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Cores SPOT ──────────────────────────────────────────────────────────────
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_SPOT = RGBColor(0x66, 0x63, 0x62)
RED_SPOT  = RGBColor(0xDE, 0x0C, 0x2F)
DARK      = RGBColor(0x3F, 0x3D, 0x3C)
LIGHT_BG  = RGBColor(0xF9, 0xF9, 0xF9)
GRAY_MID  = RGBColor(0x8C, 0x89, 0x88)
RED_DARK  = RGBColor(0xA2, 0x0A, 0x25)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank layout


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_name="Montserrat", font_size=18, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True, space_before=0, space_after=0):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, l, t, w, h,
                       font_name="Montserrat", font_size=14, bold=False,
                       color=WHITE, align=PP_ALIGN.LEFT, spacing=6):
    """lines: list of (text, bold, font_size, color) or just str"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, fs, c = line, bold, font_size, color
        else:
            txt, b, fs, c = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(fs)
        run.font.bold = b
        run.font.color.rgb = c
    return txBox


def add_circle(slide, cx, cy, r, fill_color, line_color=None, line_width=None):
    from pptx.util import Emu
    l = cx - r
    t = cy - r
    shape = slide.shapes.add_shape(9, l, t, r * 2, r * 2)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_slide(prs, bg_color=None):
    slide = prs.slides.add_slide(BLANK)
    if bg_color:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide


def add_footer(slide, text="Plataforma de Reuniões", page_num=None):
    # thin red line at bottom
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_color=RED_SPOT)
    # footer text
    label = text
    if page_num:
        label = f"{text}  •  {page_num}"
    add_text_box(slide, label, Inches(0.4), H - Inches(0.45), Inches(10), Inches(0.35),
                 font_size=8, color=GRAY_MID, align=PP_ALIGN.LEFT)


def add_section_tag(slide, text):
    """small red pill label at top left"""
    add_rect(slide, Inches(0.5), Inches(0.3), Inches(2.8), Inches(0.32), fill_color=RED_SPOT)
    add_text_box(slide, text.upper(), Inches(0.52), Inches(0.28), Inches(2.7), Inches(0.36),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=DARK)

# left red strip
add_rect(slide, 0, 0, Inches(0.25), H, fill_color=RED_SPOT)

# decorative circles (brand graphism)
add_circle(slide, W - Inches(1.8), Inches(1.2), Inches(1.8), RGBColor(0x4A, 0x47, 0x46))
add_circle(slide, W - Inches(1.1), Inches(5.8), Inches(0.9), RED_DARK)
add_circle(slide, W - Inches(3.5), H, Inches(1.2), RGBColor(0x4A, 0x47, 0x46))

# SPOT badge
add_rect(slide, Inches(0.7), Inches(0.9), Inches(1.2), Inches(0.45), fill_color=RED_SPOT)
add_text_box(slide, "SPOT", Inches(0.72), Inches(0.88), Inches(1.2), Inches(0.48),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_text_box(slide, "Plataforma de\nReuniões Inteligentes",
             Inches(0.7), Inches(1.7), Inches(8.5), Inches(2.4),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Red accent line
add_rect(slide, Inches(0.7), Inches(4.2), Inches(5.0), Inches(0.06), fill_color=RED_SPOT)

# Subtitle
add_text_box(slide, "Capture, transcreva e consulte reuniões com IA —\nsem perder nenhuma decisão, responsável ou prazo.",
             Inches(0.7), Inches(4.4), Inches(8.5), Inches(1.2),
             font_size=17, color=GRAY_MID, align=PP_ALIGN.LEFT)

# Bottom tag
add_text_box(slide, "gravação  •  transcrição  •  resumo  •  itens de ação  •  chat contextual",
             Inches(0.7), H - Inches(1.0), Inches(10), Inches(0.5),
             font_size=11, color=GRAY_MID, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O PROBLEMA
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "O Problema")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "02")

# Decorative circle background
add_circle(slide, W + Inches(0.5), Inches(-0.5), Inches(3.5), LIGHT_BG)

# Title
add_text_box(slide, "Reuniões acontecem.\nAs decisões se perdem.",
             Inches(0.5), Inches(0.85), Inches(9), Inches(1.8),
             font_size=36, bold=True, color=DARK, align=PP_ALIGN.LEFT)

# Red line
add_rect(slide, Inches(0.5), Inches(2.75), Inches(3.0), Inches(0.05), fill_color=RED_SPOT)

# 4 problem cards
cards = [
    ("📁", "Informações dispersas", "Gravações em um lugar, anotações em outro, decisões perdidas em e-mails e chats."),
    ("⏱", "Horas perdidas em atas", "Redigir atas manualmente consome tempo valioso e ainda gera documentos incompletos."),
    ("❓", "Sem rastreabilidade", "Quem ficou responsável? Qual foi o prazo combinado? Ninguém lembra ao certo."),
    ("🔍", "Difícil de consultar depois", "Voltar a uma gravação para buscar um ponto específico é trabalhoso e pouco prático."),
]

for i, (icon, title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    cl = Inches(0.5) + col * Inches(6.3)
    ct = Inches(3.1) + row * Inches(1.7)
    cw = Inches(5.9)
    ch = Inches(1.5)

    add_rect(slide, cl, ct, cw, ch, fill_color=LIGHT_BG)
    # red left accent
    add_rect(slide, cl, ct, Inches(0.07), ch, fill_color=RED_SPOT)
    add_text_box(slide, f"{icon}  {title}", cl + Inches(0.18), ct + Inches(0.12), cw - Inches(0.3), Inches(0.38),
                 font_size=13, bold=True, color=DARK)
    add_text_box(slide, body, cl + Inches(0.18), ct + Inches(0.5), cw - Inches(0.3), Inches(0.9),
                 font_size=11, color=GRAY_SPOT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — O QUE É
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=RED_SPOT)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "03")

# Decorative circles
add_circle(slide, W - Inches(1.5), Inches(1.0), Inches(2.5), RED_DARK)
add_circle(slide, W - Inches(0.5), H - Inches(0.5), Inches(1.5), RED_DARK)

add_text_box(slide, "O que é?", Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
             font_size=14, bold=True, color=RGBColor(0xFF, 0xAA, 0xAA))

add_text_box(slide, "Uma plataforma centralizada para\ncaptura, transcrição e organização de reuniões",
             Inches(0.8), Inches(1.3), Inches(10.5), Inches(2.0),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(slide, Inches(0.8), Inches(3.45), Inches(3.8), Inches(0.06), fill_color=WHITE)

add_text_box(slide,
             "Presenciais ou online, gravadas ao vivo ou enviadas como arquivo — "
             "a plataforma transcreve, resume, extrai itens de ação e permite consulta "
             "posterior por contexto, tudo em um único ambiente.",
             Inches(0.8), Inches(3.65), Inches(11.0), Inches(1.5),
             font_size=16, color=WHITE, align=PP_ALIGN.LEFT)

# 3 badges
badges = ["Capture", "Transcreva", "Consulte com IA"]
for i, b in enumerate(badges):
    bx = Inches(0.8) + i * Inches(3.6)
    add_rect(slide, bx, H - Inches(1.5), Inches(3.2), Inches(0.55),
             fill_color=RED_DARK)
    add_text_box(slide, b, bx, H - Inches(1.55), Inches(3.2), Inches(0.65),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CAPTURA FLEXÍVEL
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "Captura")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "04")

add_text_box(slide, "Grave de onde estiver,\ndo jeito que preferir",
             Inches(0.5), Inches(0.8), Inches(9), Inches(1.7),
             font_size=34, bold=True, color=DARK)
add_rect(slide, Inches(0.5), Inches(2.6), Inches(4.0), Inches(0.05), fill_color=RED_SPOT)

sources = [
    ("🎙️", "Microfone\nao Vivo",
     "Grave diretamente pelo celular ou computador durante reuniões presenciais."),
    ("📂", "Upload\nde Arquivo",
     "Envie gravações já existentes: MP3, WAV, M4A, AAC, MP4 — até 500 MB."),
    ("🖥️", "Captura\nde Reunião Online",
     "Companion desktop captura o áudio do sistema em reuniões no Teams, Zoom e Meet."),
]

for i, (icon, title, body) in enumerate(sources):
    xl = Inches(0.5) + i * Inches(4.2)
    yt = Inches(2.9)
    bw = Inches(3.9)
    bh = Inches(3.6)

    # card
    add_rect(slide, xl, yt, bw, bh, fill_color=LIGHT_BG)
    # top red accent
    add_rect(slide, xl, yt, bw, Inches(0.08), fill_color=RED_SPOT)

    # circle icon background
    add_circle(slide, xl + Inches(0.7), yt + Inches(0.9), Inches(0.55), RED_SPOT)
    add_text_box(slide, icon, xl + Inches(0.22), yt + Inches(0.4), Inches(1.0), Inches(0.9),
                 font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, title, xl + Inches(0.18), yt + Inches(1.6), bw - Inches(0.36), Inches(0.8),
                 font_size=16, bold=True, color=DARK)
    add_text_box(slide, body, xl + Inches(0.18), yt + Inches(2.4), bw - Inches(0.36), Inches(1.0),
                 font_size=12, color=GRAY_SPOT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TRANSCRIÇÃO AUTOMÁTICA
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=DARK)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "05")

add_circle(slide, Inches(1.5), H + Inches(0.5), Inches(2.5), RGBColor(0x4A, 0x47, 0x46))
add_circle(slide, W - Inches(0.8), Inches(2.0), Inches(1.8), RED_DARK)

add_text_box(slide, "TRANSCRIÇÃO", Inches(0.8), Inches(0.5), Inches(6), Inches(0.45),
             font_size=10, bold=True, color=RED_SPOT)

add_text_box(slide, "Do áudio às palavras,\nautomaticamente",
             Inches(0.8), Inches(1.0), Inches(9), Inches(2.0),
             font_size=38, bold=True, color=WHITE)

add_rect(slide, Inches(0.8), Inches(3.15), Inches(4.0), Inches(0.06), fill_color=RED_SPOT)

# Feature list
features = [
    ("🗣️", "Identificação de Falantes", "Cada participante é reconhecido e rotulado separadamente."),
    ("⏱️", "Marcação de Tempo", "Cada trecho vinculado ao momento exato da gravação."),
    ("🇧🇷", "Português Nativo", "Transcrição otimizada para o português brasileiro."),
    ("🔄", "Reprocessamento", "Falhou? Reenvie com um clique e receba nova transcrição."),
]

for i, (icon, title, desc) in enumerate(features):
    row = i % 2
    col = i // 2
    xl = Inches(0.8) + col * Inches(6.0)
    yt = Inches(3.4) + row * Inches(1.4)
    add_text_box(slide, f"{icon}  {title}", xl, yt, Inches(5.5), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)
    add_text_box(slide, desc, xl + Inches(0.05), yt + Inches(0.4), Inches(5.4), Inches(0.7),
                 font_size=11, color=GRAY_MID)
    add_rect(slide, xl, yt + Inches(1.2), Inches(5.0), Inches(0.02), fill_color=GRAY_SPOT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — IA: RESUMO, DESTAQUES E ITENS DE AÇÃO
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "Inteligência Artificial")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "06")

add_text_box(slide, "A reunião acabou.\nA IA já fez a ata.", Inches(0.5), Inches(0.8),
             Inches(9), Inches(1.7), font_size=34, bold=True, color=DARK)
add_rect(slide, Inches(0.5), Inches(2.58), Inches(3.2), Inches(0.05), fill_color=RED_SPOT)

# 3 AI output cards
ai_items = [
    ("📝", "Resumo Executivo",
     "Visão geral da reunião em 2–3 parágrafos + capítulos organizados por tema, gerados automaticamente."),
    ("⭐", "Destaques",
     "Os pontos-chave discutidos, extraídos e listados para leitura rápida — sem precisar ouvir tudo de novo."),
    ("✅", "Itens de Ação",
     "Tarefas, responsáveis e próximos passos identificados na conversa e entregues já com checkbox."),
]

for i, (icon, title, body) in enumerate(ai_items):
    xl = Inches(0.5) + i * Inches(4.2)
    yt = Inches(2.8)
    bw = Inches(3.9)
    bh = Inches(3.8)
    add_rect(slide, xl, yt, bw, bh, fill_color=LIGHT_BG)
    add_rect(slide, xl, yt, bw, Inches(0.07), fill_color=RED_SPOT)
    # number circle
    add_circle(slide, xl + Inches(0.55), yt + Inches(0.75), Inches(0.45), RED_SPOT)
    add_text_box(slide, str(i + 1), xl + Inches(0.18), yt + Inches(0.38), Inches(0.75), Inches(0.7),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"{icon} {title}", xl + Inches(0.18), yt + Inches(1.35), bw - Inches(0.3), Inches(0.5),
                 font_size=14, bold=True, color=DARK)
    add_text_box(slide, body, xl + Inches(0.18), yt + Inches(1.9), bw - Inches(0.3), Inches(1.6),
                 font_size=11, color=GRAY_SPOT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CHAT CONTEXTUAL
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=RED_SPOT)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "07")

add_circle(slide, W - Inches(1.2), H - Inches(0.8), Inches(2.2), RED_DARK)
add_circle(slide, Inches(0.3), Inches(0.3), Inches(1.0), RED_DARK)

add_text_box(slide, "CONSULTA CONTEXTUAL", Inches(0.8), Inches(0.55), Inches(6), Inches(0.4),
             font_size=10, bold=True, color=RGBColor(0xFF, 0xAA, 0xAA))
add_text_box(slide, "Pergunte sobre\na reunião — a IA responde.",
             Inches(0.8), Inches(1.0), Inches(9.5), Inches(1.8),
             font_size=38, bold=True, color=WHITE)
add_rect(slide, Inches(0.8), Inches(2.95), Inches(4.5), Inches(0.06), fill_color=WHITE)

# Mock chat bubbles
bubbles = [
    (True,  "Quem ficou responsável pelo relatório mencionado na reunião?"),
    (False, "Com base na transcrição, foi definido que Ana Silva ficaria responsável pelo relatório de vendas, com prazo até sexta-feira."),
    (True,  "Quais foram os próximos passos acordados?"),
    (False, "Foram definidos 3 itens de ação: 1) Ana entrega o relatório, 2) Carlos ajusta o cronograma e 3) Reunião de follow-up na próxima semana."),
]

by = Inches(3.2)
for is_user, text in bubbles:
    bw = Inches(7.5)
    bx = Inches(0.8) if not is_user else Inches(4.3)
    fc = WHITE if not is_user else RED_DARK
    tc = DARK if not is_user else WHITE
    add_rect(slide, bx, by, bw, Inches(0.55), fill_color=fc)
    add_text_box(slide, text, bx + Inches(0.15), by + Inches(0.06), bw - Inches(0.25), Inches(0.5),
                 font_size=10, color=tc)
    by += Inches(0.72)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROJETOS E COLABORAÇÃO
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "Organização")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "08")

add_text_box(slide, "Tudo organizado por projeto,\nacessível pela equipe.",
             Inches(0.5), Inches(0.8), Inches(9.5), Inches(1.7),
             font_size=34, bold=True, color=DARK)
add_rect(slide, Inches(0.5), Inches(2.6), Inches(3.5), Inches(0.05), fill_color=RED_SPOT)

# Left column: feature list
feats = [
    ("📁", "Projetos Temáticos",
     "Agrupe gravações por time, cliente ou iniciativa. Cada projeto tem seus próprios membros e gravações."),
    ("👥", "Papéis na Equipe",
     "Proprietários e membros com acesso controlado — quem precisa ver, consegue. Quem não deve, não acessa."),
    ("📊", "Histórico Centralizado",
     "Todas as reuniões do projeto em um só lugar, ordenadas e pesquisáveis."),
    ("🔒", "Acesso Granular",
     "Administradores têm visão total da plataforma; usuários veem apenas o que lhes é pertinente."),
]
for i, (icon, title, desc) in enumerate(feats):
    yt = Inches(2.9) + i * Inches(1.0)
    add_text_box(slide, f"{icon}  {title}", Inches(0.5), yt, Inches(5.5), Inches(0.4),
                 font_size=13, bold=True, color=DARK)
    add_text_box(slide, desc, Inches(0.5), yt + Inches(0.38), Inches(5.4), Inches(0.55),
                 font_size=11, color=GRAY_SPOT)
    if i < len(feats) - 1:
        add_rect(slide, Inches(0.5), yt + Inches(0.88), Inches(5.2), Inches(0.02), fill_color=LIGHT_BG)

# Right decorative block
add_rect(slide, Inches(7.2), Inches(2.6), Inches(5.6), Inches(4.5), fill_color=LIGHT_BG)
add_rect(slide, Inches(7.2), Inches(2.6), Inches(5.6), Inches(0.08), fill_color=RED_SPOT)

stages = [
    ("🏢  Projeto: Reuniões de Vendas", True),
    ("    📌 Kick-off Q2 — Pronto ✅", False),
    ("    📌 Sprint Review — Pronto ✅", False),
    ("    📌 Pipeline Jul/Ago — Processando ⏳", False),
    ("", False),
    ("🏢  Projeto: Onboarding Novos Colabor.", True),
    ("    📌 Apresentação Institucional — Pronto ✅", False),
    ("    📌 Treinamento Sistemas — Pronto ✅", False),
]

for i, (line, is_header) in enumerate(stages):
    yt = Inches(3.0) + i * Inches(0.45)
    add_text_box(slide, line, Inches(7.4), yt, Inches(5.2), Inches(0.42),
                 font_size=11, bold=is_header, color=DARK if is_header else GRAY_SPOT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EXPORTAÇÃO DE ATAS
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=DARK)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "09")

add_circle(slide, W - Inches(0.5), Inches(-0.5), Inches(3.0), RGBColor(0x4A, 0x47, 0x46))
add_circle(slide, Inches(0.2), H + Inches(0.2), Inches(1.8), RED_DARK)

add_text_box(slide, "EXPORTAÇÃO", Inches(0.8), Inches(0.55), Inches(5), Inches(0.4),
             font_size=10, bold=True, color=RED_SPOT)
add_text_box(slide, "Ata pronta\nem segundos.",
             Inches(0.8), Inches(1.0), Inches(9.5), Inches(1.8),
             font_size=42, bold=True, color=WHITE)
add_rect(slide, Inches(0.8), Inches(2.9), Inches(3.5), Inches(0.06), fill_color=RED_SPOT)

add_text_box(slide,
             "Exporte a ata completa com um clique: resumo, destaques, itens de ação e transcrição "
             "completa com marcação de tempo — tudo formatado e pronto para compartilhar.",
             Inches(0.8), Inches(3.1), Inches(11.0), Inches(1.0),
             font_size=14, color=GRAY_MID)

# Format cards
formats = [
    ("📄", "Markdown (.md)",
     "Formatação estruturada\ncom títulos, listas e\ncheckboxes — ideal para\nnotion, confluence, git."),
    ("📃", "Texto Plano (.txt)",
     "Compatível com qualquer\nferramenta: e-mail, word,\nnotepad — sem perda\nde informação."),
]
for i, (icon, fmt, desc) in enumerate(formats):
    xl = Inches(0.8) + i * Inches(5.0)
    yt = Inches(4.3)
    add_rect(slide, xl, yt, Inches(4.6), Inches(2.7), fill_color=RGBColor(0x4A, 0x47, 0x46))
    add_rect(slide, xl, yt, Inches(0.07), Inches(2.7), fill_color=RED_SPOT)
    add_text_box(slide, f"{icon}  {fmt}", xl + Inches(0.18), yt + Inches(0.2), Inches(4.2), Inches(0.5),
                 font_size=15, bold=True, color=WHITE)
    add_text_box(slide, desc, xl + Inches(0.18), yt + Inches(0.75), Inches(4.2), Inches(1.7),
                 font_size=12, color=GRAY_MID)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BUSCA E ACESSO POSTERIOR
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "Busca & Acesso")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "10")

add_text_box(slide, "Encontre qualquer informação\nde qualquer reunião passada.",
             Inches(0.5), Inches(0.8), Inches(12.5), Inches(1.7),
             font_size=34, bold=True, color=DARK)
add_rect(slide, Inches(0.5), Inches(2.6), Inches(5.0), Inches(0.05), fill_color=RED_SPOT)

# Search feature highlight
add_rect(slide, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.0), fill_color=LIGHT_BG)
add_rect(slide, Inches(0.5), Inches(2.8), Inches(0.07), Inches(1.0), fill_color=RED_SPOT)
add_text_box(slide, "🔍  Busca inteligente por título, resumo, transcrição completa e notas — "
             "uma única caixa de busca encontra em qualquer gravação.",
             Inches(0.7), Inches(2.9), Inches(11.8), Inches(0.75),
             font_size=14, bold=False, color=DARK)

# Filter chips
filters = ["Por projeto", "Por status", "Por data", "Sem projeto", "Em processamento"]
fx = Inches(0.5)
fy = Inches(4.1)
for f in filters:
    add_rect(slide, fx, fy, Inches(len(f) * 0.11 + 0.5), Inches(0.38),
             fill_color=LIGHT_BG, line_color=GRAY_MID, line_width=Pt(0.5))
    add_text_box(slide, f, fx + Inches(0.08), fy + Inches(0.04), Inches(len(f) * 0.11 + 0.35), Inches(0.33),
                 font_size=11, color=GRAY_SPOT)
    fx += Inches(len(f) * 0.11 + 0.65)

# Access stages
access_items = [
    ("🟢", "Pronto",        "Acesso imediato a resumo, chat, exportação e transcrição."),
    ("🟡", "Em Processamento", "Acompanhe o status em tempo real na biblioteca."),
    ("🔴", "Falhou",        "Reprocesse com um clique — sem perder a gravação original."),
]
for i, (dot, status, desc) in enumerate(access_items):
    yt = Inches(4.7) + i * Inches(0.85)
    add_text_box(slide, f"{dot}  {status}", Inches(0.5), yt, Inches(4.0), Inches(0.4),
                 font_size=13, bold=True, color=DARK)
    add_text_box(slide, desc, Inches(0.5), yt + Inches(0.38), Inches(12.0), Inches(0.38),
                 font_size=11, color=GRAY_SPOT)
    if i < 2:
        add_rect(slide, Inches(0.5), yt + Inches(0.76), Inches(12.2), Inches(0.02), fill_color=LIGHT_BG)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MULTIPLATAFORMA
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=WHITE)
add_section_tag(slide, "Plataformas")
add_footer(slide, "Plataforma de Reuniões Inteligentes", "11")

add_text_box(slide, "Onde você trabalha,\na plataforma está.",
             Inches(0.5), Inches(0.8), Inches(9.5), Inches(1.7),
             font_size=34, bold=True, color=DARK)
add_rect(slide, Inches(0.5), Inches(2.6), Inches(3.0), Inches(0.05), fill_color=RED_SPOT)

platforms = [
    ("🌐", "Web", "Acesse pelo navegador, sem instalar nada. Funciona em qualquer computador."),
    ("📱", "Mobile\niOS & Android", "Grave ao vivo com o microfone do celular. Acesse e consulte de qualquer lugar."),
    ("🖥️", "Desktop\nWindows & macOS", "Companion que captura o áudio das reuniões online diretamente do sistema."),
]

for i, (icon, plat, desc) in enumerate(platforms):
    xl = Inches(0.5) + i * Inches(4.2)
    yt = Inches(2.9)
    bw = Inches(3.9)
    bh = Inches(3.8)
    add_rect(slide, xl, yt, bw, bh, fill_color=LIGHT_BG)
    add_rect(slide, xl, yt, bw, Inches(0.07), fill_color=RED_SPOT)

    # big icon
    add_text_box(slide, icon, xl, yt + Inches(0.4), bw, Inches(1.0),
                 font_size=36, color=RED_SPOT, align=PP_ALIGN.CENTER)
    add_text_box(slide, plat, xl + Inches(0.15), yt + Inches(1.5), bw - Inches(0.3), Inches(0.7),
                 font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, xl + Inches(0.18), yt + Inches(2.25), bw - Inches(0.36), Inches(1.3),
                 font_size=11, color=GRAY_SPOT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — VISÃO GERAL DO FLUXO
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=DARK)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "12")

add_circle(slide, W - Inches(1.0), Inches(1.0), Inches(2.2), RGBColor(0x4A, 0x47, 0x46))

add_text_box(slide, "Do áudio à ata: tudo automático.", Inches(0.8), Inches(0.6),
             Inches(11), Inches(0.8), font_size=32, bold=True, color=WHITE)
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.0), Inches(0.05), fill_color=RED_SPOT)

steps = [
    ("1", "Grave ou\nenviAu",     "Microfone, upload ou\ncaptura de reunião online."),
    ("2", "Transcrição\nAutomática", "Identificação de\nfalantes e timestamps."),
    ("3", "IA Processa",          "Resumo, destaques,\ntags e itens de ação."),
    ("4", "Consulte\ne Exporte",  "Chat contextual,\nbusca e exportação."),
]

step_w = Inches(2.8)
step_h = Inches(3.8)
step_gap = Inches(0.5)
start_x = Inches(0.8)

for i, (num, title, desc) in enumerate(steps):
    xl = start_x + i * (step_w + step_gap)
    yt = Inches(2.0)

    add_rect(slide, xl, yt, step_w, step_h, fill_color=RGBColor(0x4A, 0x47, 0x46))
    add_rect(slide, xl, yt, step_w, Inches(0.07), fill_color=RED_SPOT)

    # number circle
    add_circle(slide, xl + step_w / 2, yt + Inches(0.9), Inches(0.55), RED_SPOT)
    add_text_box(slide, num, xl + step_w / 2 - Inches(0.3), yt + Inches(0.5), Inches(0.6), Inches(0.75),
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, title, xl + Inches(0.15), yt + Inches(1.75), step_w - Inches(0.3), Inches(0.8),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, xl + Inches(0.15), yt + Inches(2.55), step_w - Inches(0.3), Inches(1.0),
                 font_size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)

    # arrow between
    if i < len(steps) - 1:
        ax = xl + step_w + Inches(0.12)
        add_text_box(slide, "→", ax, yt + Inches(1.5), Inches(0.35), Inches(0.5),
                     font_size=18, bold=True, color=RED_SPOT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ENCERRAMENTO
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, bg_color=RED_SPOT)
add_footer(slide, "Plataforma de Reuniões Inteligentes", "13")

add_circle(slide, W - Inches(1.5), Inches(1.2), Inches(2.5), RED_DARK)
add_circle(slide, Inches(0.5), H - Inches(0.5), Inches(1.5), RED_DARK)
add_circle(slide, W / 2, H + Inches(0.8), Inches(2.0), RED_DARK)

# SPOT badge
add_rect(slide, Inches(0.8), Inches(0.7), Inches(1.2), Inches(0.44), fill_color=RED_DARK)
add_text_box(slide, "SPOT", Inches(0.82), Inches(0.68), Inches(1.2), Inches(0.48),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Reuniões produtivas.\nDecisões documentadas.\nEquipes alinhadas.",
             Inches(0.8), Inches(1.4), Inches(10), Inches(2.6),
             font_size=42, bold=True, color=WHITE)

add_rect(slide, Inches(0.8), Inches(4.15), Inches(4.5), Inches(0.06), fill_color=WHITE)

add_text_box(slide,
             "Com a plataforma de reuniões SPOT, sua equipe para de perder tempo procurando "
             "informações e começa a agir com base em registros claros, acessíveis e inteligentes.",
             Inches(0.8), Inches(4.35), Inches(11.0), Inches(1.2),
             font_size=15, color=WHITE)

# bottom tag line
add_text_box(slide,
             "gravação  •  transcrição automática  •  resumo com IA  •  itens de ação  •  chat contextual  •  exportação de atas",
             Inches(0.8), H - Inches(1.1), Inches(11.5), Inches(0.5),
             font_size=10, color=RGBColor(0xFF, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "C:/vscode_projects/Plaude_like/Apresentacao_Plataforma_Reunioes.pptx"
prs.save(out_path)
print(f"Salvo em: {out_path}")
